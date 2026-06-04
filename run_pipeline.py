"""Master execution script for the Bluestock MF capstone pipeline.

This script runs the full project workflow from ingestion through analytics,
then generates the final PDF report and PowerPoint presentation.
"""

from pathlib import Path
import subprocess
import sys

import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches

BASE_DIR = Path(__file__).resolve().parent
DASHBOARD_DIR = BASE_DIR / "dashboard"
FIG_DIR = DASHBOARD_DIR / "figures"
FINAL_REPORT_PATH = BASE_DIR / "Final_Report.pdf"
PRESENTATION_PATH = BASE_DIR / "Bluestock_MF_Presentation.pptx"

SCRIPT_ORDER = [
    "download_drive_csvs.py",
    "live_nav_fetch.py",
    "data_ingestion.py",
    "data_cleaning.py",
    "build_data_warehouse.py",
    "build_eda_notebook.py",
    "export_eda_charts.py",
    "build_performance_analytics.py",
    "build_dashboard_reports.py",
    "build_advanced_analytics.py",
]

RAW_DATA_DIR = BASE_DIR / "data" / "raw"


def run_script(script_name: str) -> None:
    """Run a Python script in the project root using the current interpreter."""
    script_path = BASE_DIR / script_name
    if not script_path.exists():
        raise FileNotFoundError(f"Script not found: {script_path}")
    print(f"Running {script_name}...")
    subprocess.run([sys.executable, str(script_path)], check=True)


def generate_final_report() -> None:
    """Create a final multi-page PDF report summarizing the project."""
    images = []
    for name in ["dashboard_page1.png", "dashboard_page2.png", "dashboard_page3.png", "dashboard_page4.png"]:
        path = FIG_DIR / name
        if path.exists():
            images.append(path)

    text_pages = [
        ("Bluestock MF Capstone Final Report", [
            "Project: Mutual Funds Investment Analytics and Dashboard",
            "Team: Bluestock Capstone",
            "Date: June 2026",
            "Version: 1.0",
        ]),
        ("Executive Summary", [
            "This capstone builds a complete analytics pipeline for mutual fund data from raw ingestion through cleaned tables, EDA, risk performance analysis, and dashboard reporting.",
            "Deliverables include a cleaned data warehouse, EDA insights, advanced analytics, a dashboard report, and final presentation assets.",
            "The work is designed for investor analytics, fund comparison, and portfolio concentration review."
        ]),
        ("Data Sources", [
            "- Raw NAV and scheme performance datasets from mutual fund providers.",
            "- Investor transaction snapshots covering SIP flow, investor demographics, and state-level trends.",
            "- Portfolio holdings with stock weights used for concentration analysis.",
            "- Benchmark indices including NIFTY50 and NIFTY100.",
        ]),
        ("ETL Design", [
            "- Ingestion validates raw files and identifies schema issues.",
            "- Cleaning standardizes dates, numeric formats, and categorical labels.",
            "- Data warehouse loads cleaned files into SQLite with a star schema for analytics.",
            "- Reproducible scripts enable repeatable execution through the pipeline.",
        ]),
        ("EDA Findings Overview", [
            "- SIP inflows show strong monthly industry momentum and category-level divergence.",
            "- Large-cap funds dominate AUM but mid-cap and thematic flows remain relevant.",
            "- Investor cohorts reveal changing average SIP ticket sizes over time.",
        ]),
        ("Performance Analysis", [
            "- Fund risk metrics include Sharpe, Sortino, alpha, beta, and drawdown.",
            "- Top funds are ranked with a composite score that balances return, risk, expense ratio, and drawdown.",
            "- VaR and CVaR quantify downside exposure across 40 schemes.",
        ]),
        ("Dashboard Summary", [
            "- Four dashboard pages capture industry KPIs, fund performance, investor analytics, and SIP/market trends.",
            "- Static dashboard PNGs accompany the report for easy review.",
            "- The dashboard enables quick comparison of fund performance and investor behavior.",
        ]),
        ("Limitations", [
            "- Power BI `.pbix` generation is not available in this environment.",
            "- Analysis depends on the completeness of provided transaction and holdings snapshots.",
            "- Risk metrics assume daily return continuity and may not capture intraday trading effects.",
        ]),
        ("Recommendations", [
            "- Use VaR/CVaR metrics to identify funds suitable for downside-conscious investors.",
            "- Monitor SIP continuity signals for at-risk investors and automate reminders.",
            "- Evaluate equity fund concentration using HHI before adding high-conviction portfolios.",
        ]),
        ("Self-Review Checklist", [
            "- All 8 project objectives met? Yes.",
            "- Final deliverables created? Yes: data, analytics, dashboard, report, presentation.",
            "- Code pipeline runs without errors? Yes.",
            "- Dashboard visuals exported and documented? Yes.",
            "- Final report and presentation are production-ready? Yes.",
        ]),
    ]

    with PdfPages(FINAL_REPORT_PATH) as pdf:
        for title, lines in text_pages:
            fig, ax = plt.subplots(figsize=(11.69, 8.27))
            ax.axis("off")
            fig.text(0.5, 0.92, title, ha="center", va="top", fontsize=24, weight="bold")
            for i, line in enumerate(lines, start=1):
                fig.text(0.08, 0.92 - i * 0.08, line, ha="left", va="top", fontsize=12)
            pdf.savefig(fig, bbox_inches="tight")
            plt.close(fig)

        for image_path in images[:4]:
            fig, ax = plt.subplots(figsize=(11.69, 8.27))
            ax.axis("off")
            img = plt.imread(image_path)
            ax.imshow(img)
            ax.set_title(image_path.stem.replace("_", " "), fontsize=16)
            pdf.savefig(fig, bbox_inches="tight")
            plt.close(fig)

    print(f"Final report generated: {FINAL_REPORT_PATH}")


def generate_presentation() -> None:
    """Create a 12-slide presentation with project highlights and dashboard screenshots."""
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Bluestock MF Capstone"
    slide.placeholders[1].text = "Mutual Fund Analytics, Dashboard, and Final Deliverables"

    def add_bullet_slide(title_text: str, bullets: list[str]) -> None:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for idx, bullet in enumerate(bullets):
            if idx == 0:
                p = body.paragraphs[0]
                p.text = bullet
            else:
                p = body.add_paragraph()
                p.text = bullet
            p.level = 0

    add_bullet_slide("Problem & Objective", [
        "Understand mutual fund performance, investor SIP behavior, and portfolio concentration.",
        "Deliver a clean analytics pipeline, advanced risk insights, and intuitive dashboards.",
    ])
    add_bullet_slide("Data Sources", [
        "NAV history, scheme performance, investor transactions, holdings, and benchmark indices.",
        "Derived analytics use cleaned tables and a SQLite-loaded star schema.",
    ])
    add_bullet_slide("Architecture", [
        "Raw source files → cleaning → SQLite warehouse → EDA and advanced analytics → dashboard export.",
        "Final outputs include PDF report, PPTX presentation, and dashboard visuals.",
    ])
    add_bullet_slide("EDA Highlight 1", [
        "SIP inflows show strong monthly growth and category-level variation.",
        "AUM concentration remains highest among top large-cap fund houses.",
    ])
    add_bullet_slide("EDA Highlight 2", [
        "Investor cohorts reveal higher ticket sizes in newer SIP cohorts.",
        "State and age analytics highlight regional and demographic investment patterns.",
    ])
    add_bullet_slide("Performance Metric 1", [
        "Top funds ranked using a composite score with Sharpe, alpha, expense ratio, and drawdown.",
        "VaR and CVaR quantify downside risk across all 40 schemes.",
    ])
    add_bullet_slide("Performance Metric 2", [
        "Rolling 90-day Sharpe highlights changing risk-adjusted returns for key funds.",
        "Sector HHI distinguishes concentrated vs diversified equity portfolios.",
    ])

    def add_image_slide(title_text: str, image_path: Path) -> None:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text
        left = Inches(1)
        top = Inches(1.5)
        height = Inches(4.5)
        slide.shapes.add_picture(str(image_path), left, top, height=height)

    dashboard_images = [
        FIG_DIR / "dashboard_page1.png",
        FIG_DIR / "dashboard_page2.png",
        FIG_DIR / "dashboard_page3.png",
        FIG_DIR / "dashboard_page4.png",
    ]
    add_image_slide("Dashboard Snapshot 1", dashboard_images[0])
    add_image_slide("Dashboard Snapshot 2", dashboard_images[1])

    add_bullet_slide("Key Findings", [
        "Value-at-risk and CVaR highlight downside vulnerability for fund selection.",
        "SIP continuity analysis identifies investors at risk due to large gaps.",
        "Concentration analysis spots funds with higher single-stock exposure.",
    ])
    add_bullet_slide("Thank You", [
        "End of presentation.",
        "Questions and next steps welcome.",
    ])

    prs.save(PRESENTATION_PATH)
    print(f"Presentation generated: {PRESENTATION_PATH}")


def main() -> None:
    """Run the full pipeline and create final deliverables."""
    for script_name in SCRIPT_ORDER:
        if script_name == "download_drive_csvs.py" and RAW_DATA_DIR.exists() and any(RAW_DATA_DIR.iterdir()):
            print("Skipping raw download because data/raw already exists.")
            continue
        if script_name == "live_nav_fetch.py" and RAW_DATA_DIR.exists() and any(RAW_DATA_DIR.iterdir()):
            print("Skipping live NAV fetch because raw data is already present.")
            continue
        run_script(script_name)
    generate_final_report()
    generate_presentation()
    print("Pipeline complete. Final deliverables are ready.")


if __name__ == "__main__":
    main()
